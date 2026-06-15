"""
PPTReviewerAgent - 演示文稿审核专家

职责：审核演示文稿作品（PPT_TRACK）

官方评分维度（校教发〔2024〕77号文件）：
- 创意（内容创意、视觉设计创新）: 0-20分
- 视觉效果（排版、色彩、图文比例）: 0-20分
- 内容呈现（逻辑结构、信息密度）: 0-20分
- 原创性（原创元素使用）: 0-20分

总分：100分（AI评分 × weight + 评委评分 × weight）

硬性要求检查：
- 页数：至少12页
- 比例：16:9
- 格式：PPTX，无密码保护
- 原创性（不得抄袭）
- 内容健康积极

工具集成（未来实现）：
- PPTStructureAnalyzerTool：幻灯片结构分析
- PPTMetadataTool：元数据检查
- DesignQualityTool：设计质量评估
- OriginalityDetectorTool：原创性检测
"""

import os
import json
import logging
from typing import Dict, Any, Optional, List
from pydantic import BaseModel, Field
from langchain_openai import ChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage
from app.utils.llm_utils import extract_json_from_llm_output

logger = logging.getLogger(__name__)


class PPTReviewOutput(BaseModel):
    """PPT评审结构化输出（官方评分维度）"""

    # 总分（100分）
    overall_score: float = Field(description="总体评分（0-100分）", ge=0, le=100)

    # 官方评分维度（各20分，5维度合计100分）
    creativity_score: float = Field(description="创意评分（0-20分）", ge=0, le=20)
    visual_effect_score: float = Field(description="视觉效果评分（0-20分）", ge=0, le=20)
    content_presentation_score: float = Field(description="内容呈现评分（0-20分）", ge=0, le=20)
    originality_score: float = Field(description="原创性评分（0-20分）", ge=0, le=20)
    documentation_score: float = Field(description="文档完整性评分（0-20分）", ge=0, le=20)

    # 硬性要求合规性
    compliance_check: Dict[str, Any] = Field(description="硬性要求合规性检查结果")

    # 评审意见
    review_summary: str = Field(description="评审总结（200-500字）")
    strengths: List[str] = Field(description="作品亮点（3-5条）")
    weaknesses: List[str] = Field(description="不足之处（2-3条）")
    improvement_suggestions: List[str] = Field(description="改进建议（3-5条）")

    # PPT元数据
    ppt_metadata: Dict[str, Any] = Field(description="PPT元数据")


class PPTReviewerAgent:
    """
    PPT审核专家Agent

    官方评分维度：
    - 创意（内容创意、视觉设计创新）
    - 视觉效果（排版、色彩、图文比例）
    - 内容呈现（逻辑结构、信息密度）
    - 原创性（原创元素使用）

    硬性要求检查：
    - 页数：至少12页
    - 比例：16:9
    - 格式：PPTX，无密码保护
    """

    def __init__(
        self,
        model_name: str = "deepseek-chat",
        api_key: Optional[str] = None,
        base_url: Optional[str] = None
    ):
        """
        初始化PPT审核Agent

        Args:
            model_name: DeepSeek模型名称
            api_key: DeepSeek API密钥
            base_url: DeepSeek API地址
        """
        # 配置DeepSeek API
        self.api_key = api_key or os.getenv("DEEPSEEK_API_KEY")
        self.base_url = base_url or os.getenv("DEEPSEEK_BASE_URL", "https://api.deepseek.com")

        if not self.api_key:
            raise ValueError("DeepSeek API密钥未设置")

        # 初始化LLM
        self.llm = ChatOpenAI(
            model=model_name,
            api_key=self.api_key,
            base_url=self.base_url,
            temperature=0.3
        )

        logger.info("PPTReviewerAgent初始化完成")

    def review_ppt(
        self,
        ppt_path: str,
        work_description: str,
        readme_content: Optional[str] = None,
        additional_files: Optional[List[str]] = None
    ) -> PPTReviewOutput:
        """
        完整PPT审核流程

        Args:
            ppt_path: PPT文件路径
            work_description: 作品说明文档内容
            readme_content: README文件内容（可选）
            additional_files: 其他附加文件路径列表

        Returns:
            PPTReviewOutput: 结构化评审报告
        """
        logger.info(f"开始PPT审核：{ppt_path}")

        # ========== 步骤1：PPT元数据提取 ==========
        logger.info("步骤1：PPT元数据提取...")
        ppt_metadata = self._extract_ppt_metadata(ppt_path)

        # ========== 步骤2：硬性要求检查 ==========
        logger.info("步骤2：硬性要求检查...")
        compliance_check = self._check_compliance(ppt_metadata, readme_content)

        # ========== 步骤3：PPT质量分析（增强版） ==========
        logger.info("步骤3：PPT质量分析...")
        ppt_quality_evidence = self._analyze_ppt_quality(ppt_metadata)
        visual_effect_score = ppt_quality_evidence.get("visual_score_hint", 8)

        # ========== 步骤4：DeepSeek LLM评审推理 ==========
        logger.info("步骤4：DeepSeek LLM评审推理...")

        # 构建评审prompt
        prompt = self._build_review_prompt(
            ppt_metadata=ppt_metadata,
            compliance_check=compliance_check,
            work_description=work_description,
            readme_content=readme_content,
            visual_effect_score=visual_effect_score,
            ppt_quality_evidence=ppt_quality_evidence,
        )

        # 调用LLM
        response = self.llm.invoke([
            SystemMessage(content=self._get_system_prompt()),
            HumanMessage(content=prompt)
        ])

        # 解析LLM输出
        review_result = self._parse_llm_output(
            llm_output=response.content,
            ppt_metadata=ppt_metadata,
            compliance_check=compliance_check,
            visual_effect_score=visual_effect_score
        )

        logger.info(f"PPT审核完成，总分：{review_result.overall_score}")

        return review_result

    def _extract_ppt_metadata(self, ppt_path: str) -> Dict[str, Any]:
        """
        提取PPT元数据（基础版，使用python-pptx）

        Args:
            ppt_path: PPT文件路径

        Returns:
            dict: PPT元数据
        """
        metadata = {
            "ppt_path": ppt_path,
            "file_size": 0,
            "slide_count": 0,
            "format": "Unknown",
            "has_password": False,
            "aspect_ratio": "Unknown",
            "slide_dimensions": None,
            "slide_contents": []  # 新增：提取每页文本内容
        }

        try:
            # 检查文件是否存在
            if not os.path.exists(ppt_path):
                logger.error(f"PPT文件不存在：{ppt_path}")
                return metadata

            # 文件大小
            metadata["file_size"] = os.path.getsize(ppt_path)

            # 格式检查（文件扩展名）
            file_ext = os.path.splitext(ppt_path)[1].lower()
            if file_ext == ".pptx":
                metadata["format"] = "PPTX"
            elif file_ext == ".ppt":
                metadata["format"] = "PPT"
            else:
                metadata["format"] = "Unknown"

            # 尝试使用python-pptx提取详细信息
            try:
                from pptx import Presentation

                prs = Presentation(ppt_path)

                # 幻灯片数量
                metadata["slide_count"] = len(prs.slides)

                # 幻灯片尺寸（判断16:9比例）
                slide_width = prs.slide_width
                slide_height = prs.slide_height

                # EMUs (English Metric Units) to inches conversion
                # 1 inch = 914400 EMUs
                width_inches = slide_width / 914400
                height_inches = slide_height / 914400

                # 计算比例
                if height_inches > 0:
                    ratio = width_inches / height_inches
                    if abs(ratio - 16/9) < 0.1:  # 允许10%误差
                        metadata["aspect_ratio"] = "16:9"
                    elif abs(ratio - 4/3) < 0.1:
                        metadata["aspect_ratio"] = "4:3"
                    else:
                        metadata["aspect_ratio"] = f"{ratio:.2f}:1"

                metadata["slide_dimensions"] = {
                    "width_inches": round(width_inches, 2),
                    "height_inches": round(height_inches, 2),
                    "width_emus": slide_width,
                    "height_emus": slide_height
                }

                # 提取每页幻灯片的文本内容
                slide_contents = []
                for idx, slide in enumerate(prs.slides, start=1):
                    slide_text = {
                        "slide_number": idx,
                        "title": "",
                        "content": []
                    }

                    # 遍历所有shape提取文本
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text = shape.text.strip()
                            # 判断是否为标题（通常第一个文本框或大字体）
                            if not slide_text["title"] and len(text) < 100:
                                slide_text["title"] = text
                            else:
                                slide_text["content"].append(text)

                    # 如果该页有内容，添加到列表
                    if slide_text["title"] or slide_text["content"]:
                        slide_contents.append(slide_text)

                metadata["slide_contents"] = slide_contents
                logger.info(f"成功提取{len(slide_contents)}页幻灯片文本内容")

                # 尝试检测密码保护（通过尝试打开判断）
                # 如果成功打开，说明无密码保护
                metadata["has_password"] = False

            except ImportError:
                logger.warning("python-pptx未安装，无法提取详细元数据")
                # 基础检查：至少文件存在且格式正确
                metadata["slide_count"] = 0  # 未知

            except Exception as e:
                logger.error(f"PPT元数据提取失败：{e}")
                # 可能是密码保护导致无法打开
                if "password" in str(e).lower() or "encrypted" in str(e).lower():
                    metadata["has_password"] = True

        except Exception as e:
            logger.error(f"PPT文件读取失败：{e}")

        return metadata

    def _check_compliance(
        self,
        ppt_metadata: Dict[str, Any],
        readme_content: Optional[str]
    ) -> Dict[str, Any]:
        """
        检查硬性要求合规性

        检查项：
        1. 页数：至少12页
        2. 比例：16:9
        3. 格式：PPTX，无密码保护
        4. 文件大小合理性

        Args:
            ppt_metadata: PPT元数据
            readme_content: README文件内容

        Returns:
            dict: 合规性检查结果
        """
        compliance = {}

        # 1. 页数要求（至少12页）
        compliance["slide_count_valid"] = ppt_metadata["slide_count"] >= 12
        compliance["slide_count_message"] = (
            f"幻灯片数量{ppt_metadata['slide_count']}页，符合要求（≥12页）"
            if compliance["slide_count_valid"]
            else f"幻灯片数量{ppt_metadata['slide_count']}页，不符合要求（应≥12页）"
        )

        # 2. 比例要求（16:9）
        compliance["ratio_valid"] = ppt_metadata["aspect_ratio"] == "16:9"
        compliance["ratio_message"] = (
            f"幻灯片比例{ppt_metadata['aspect_ratio']},符合要求"
            if compliance["ratio_valid"]
            else f"幻灯片比例{ppt_metadata['aspect_ratio']},不符合要求（应为16:9）"
        )

        # 3. 格式要求（PPTX）
        compliance["format_valid"] = ppt_metadata["format"] == "PPTX"
        compliance["format_message"] = (
            f"文件格式{ppt_metadata['format']},符合要求"
            if compliance["format_valid"]
            else f"文件格式{ppt_metadata['format']},不符合要求（应为PPTX）"
        )

        # 4. 密码保护检查（无密码）
        compliance["password_valid"] = not ppt_metadata["has_password"]
        compliance["password_message"] = (
            "文件无密码保护，符合要求"
            if compliance["password_valid"]
            else "文件有密码保护，不符合要求"
        )

        # 5. 文件大小合理性（≤300MB）
        file_size_mb = ppt_metadata["file_size"] / (1024 * 1024)
        compliance["size_valid"] = file_size_mb <= 300
        compliance["size_message"] = (
            f"文件大小{file_size_mb:.2f}MB，符合要求（≤300MB）"
            if compliance["size_valid"]
            else f"文件大小{file_size_mb:.2f}MB，不符合要求（应≤300MB）"
        )

        # 6. 说明文档完整性
        compliance["readme_exists"] = readme_content is not None and len(readme_content) > 100
        compliance["readme_message"] = (
            "说明文档完整，符合要求"
            if compliance["readme_exists"]
            else "说明文档不完整或不存在"
        )

        # 7. 总体合规性
        compliance["all_valid"] = all([
            compliance["slide_count_valid"],
            compliance["ratio_valid"],
            compliance["format_valid"],
            compliance["password_valid"],
            compliance["size_valid"],
            compliance["readme_exists"]
        ])

        compliance["overall_message"] = (
            "PPT作品符合所有硬性要求"
            if compliance["all_valid"]
            else "PPT作品不符合部分硬性要求"
        )

        return compliance

    def _analyze_ppt_quality(self, ppt_metadata: Dict[str, Any]) -> Dict[str, Any]:
        """
        分析PPT质量（增强版），返回结构化分析结果

        Args:
            ppt_metadata: PPT元数据

        Returns:
            dict: PPT质量分析结果（含各维度指标和建议分数参考）
        """
        evidence = {
            "slide_count": ppt_metadata.get("slide_count", 0),
            "has_chart": False,
            "avg_text_per_slide": 0,
            "text_density_level": "未知",
            "image_slide_count": 0,
            "image_ratio": 0.0,
            "has_structure": False,
            "font_count": 0,
            "font_consistency": "未知",
            "visual_score_hint": 8,
            "content_score_hint": 8,
        }

        try:
            from pptx import Presentation

            ppt_path = ppt_metadata.get("ppt_path", "")
            if not os.path.exists(ppt_path):
                return evidence

            prs = Presentation(ppt_path)
            slides = prs.slides
            total_text_chars = 0
            text_slides = 0
            image_slides = 0
            chart_count = 0
            fonts = set()
            has_toc = False

            for slide in slides:
                slide_text_len = 0
                has_image = False

                for shape in slide.shapes:
                    # 文本统计
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_len += len(shape.text.strip())

                    # 图片检测
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        has_image = True

                    # 图表检测
                    if shape.has_chart:
                        chart_count += 1

                    # 字体收集
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name:
                                    fonts.add(run.font.name)

                total_text_chars += slide_text_len
                if slide_text_len > 0:
                    text_slides += 1
                if has_image:
                    image_slides += 1

                # 目录页检测（简单启发式：标题含"目录"/"contents"/"大纲"）
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_lower = shape.text.strip().lower()
                        if any(kw in text_lower for kw in ["目录", "contents", "大纲", "提纲", "概览"]):
                            if len(text_lower) < 30:
                                has_toc = True

            slide_count = len(slides)

            # 图文比
            evidence["image_slide_count"] = image_slides
            if slide_count > 0:
                evidence["image_ratio"] = round(image_slides / slide_count, 2)

            # 平均每页文字量
            if text_slides > 0:
                evidence["avg_text_per_slide"] = round(total_text_chars / text_slides)

            # 文字密度等级
            avg = evidence["avg_text_per_slide"]
            if avg > 300:
                evidence["text_density_level"] = "偏高（文字较多）"
            elif avg > 100:
                evidence["text_density_level"] = "适中"
            elif avg > 0:
                evidence["text_density_level"] = "偏低（文字较少）"
            else:
                evidence["text_density_level"] = "无文字"

            # 图表
            evidence["has_chart"] = chart_count > 0

            # 结构
            evidence["has_structure"] = has_toc or slide_count >= 15

            # 字体一致性
            evidence["font_count"] = len(fonts)
            if len(fonts) <= 2:
                evidence["font_consistency"] = "统一"
            elif len(fonts) <= 4:
                evidence["font_consistency"] = "基本统一"
            else:
                evidence["font_consistency"] = "较杂乱"

            # 视觉效果建议分数
            vis_score = 6
            if 0.3 <= evidence["image_ratio"] <= 0.8:
                vis_score += 3
            elif evidence["image_ratio"] > 0:
                vis_score += 1
            if evidence["has_chart"]:
                vis_score += 2
            if evidence["font_consistency"] == "统一":
                vis_score += 2
            elif evidence["font_consistency"] == "基本统一":
                vis_score += 1
            if slide_count >= 12:
                vis_score += 2
            evidence["visual_score_hint"] = min(vis_score, 20)

            # 内容呈现建议分数
            content_score = 6
            if evidence["text_density_level"] == "适中":
                content_score += 4
            elif evidence["text_density_level"] in ("偏高（文字较多）", "偏低（文字较少）"):
                content_score += 2
            if evidence["has_structure"]:
                content_score += 3
            if slide_count >= 15:
                content_score += 2
            elif slide_count >= 12:
                content_score += 1
            if evidence["avg_text_per_slide"] > 0:
                content_score += 2
            evidence["content_score_hint"] = min(content_score, 20)

        except ImportError:
            logger.warning("python-pptx未安装，跳过增强质量分析")
        except Exception as e:
            logger.warning(f"PPT增强质量分析失败：{e}")

        return evidence

    def _get_system_prompt(self) -> str:
        """系统提示词：定义PPT评审角色和标准（非计算机专业友好版）"""
        return """你是一位经验丰富的演示文稿评审专家，负责评审非计算机专业大学生的演示文稿作品。

重要背景：参赛者为非计算机专业学生，设计经验和技能参差不齐。请以鼓励性评价为主，在指出不足的同时肯定其努力和亮点。评分应体现公平性，不应以专业设计师的标准要求参赛者。

【评分维度】（5个维度各20分，总分100分）：

1. 创意性（0-20分）：
   16-20分：主题角度新颖独特，切入方式令人印象深刻；页面设计有创意（独特的配色方案、布局或动画效果）；内容呈现方式有巧思（如用故事线串联、互动环节等）
   11-15分：主题明确有一定新意；页面设计有用心之处；使用了动画或过渡效果增加表现力
   6-10分：主题常见但表述清楚；页面设计中规中矩
   0-5分：主题模糊，缺乏创意

2. 视觉效果（0-20分）：
   16-20分：整体风格统一美观；图文搭配合理，信息密度适中；配色协调，有层次感；使用了合适的图表/图片增强表达
   11-15分：整体较为整洁；有基本的图文搭配；配色基本协调，页面布局有一定逻辑
   6-10分：能看出排版意图但不够精细；图文比例失调或配色不协调
   0-5分：排版混乱，文字堆砌，无视觉设计

3. 内容呈现（0-20分）：
   16-20分：逻辑清晰完整（有开头-主体-结论结构）；信息密度适中，重点突出；文字表达准确流畅；内容有深度，有数据或案例支撑
   11-15分：有基本逻辑结构；内容较为充实但部分页面信息过多或过少；文字表达基本通顺
   6-10分：结构不够清晰；内容较为单薄或文字过多；有少量错别字或表述不当
   0-5分：内容混乱，无逻辑可言；大量错别字或语病

4. 原创性（0-20分）：
   16-20分：内容为原创，有独特的观点或数据；设计风格独特，非直接使用网络模板；包含自制的图表、图片或素材
   11-15分：主体内容为原创，部分素材来自网络但已注明来源；在模板基础上做了明显的个性化修改
   6-10分：内容基本原创但设计依赖模板；素材来源标注不完整
   0-5分：内容大量复制，未注明来源；直接使用未修改的网络模板

5. 文档完整性（0-20分）：
   16-20分：有完整README（项目介绍、设计思路、演示逻辑）；有设计说明和素材来源说明；文档内容详实、格式规范
   11-15分：有README但内容不够详实；有部分设计说明；文档基本覆盖主题但缺少细节
   6-10分：仅有简单的README或说明；文档内容简短，缺少关键信息
   0-5分：无任何说明文档，或文档内容与作品严重不符

【硬性要求检查】：
- 页数：至少12页
- 比例：16:9
- 格式：PPTX，无密码保护
- 原创性（不得抄袭）
- 内容健康积极

评审原则：
1. 参赛者为非计算机专业学生，请以鼓励性评价为主
2. 在指出不足的同时肯定其努力和亮点
3. 硬性要求不合规项需指出，但语气要温和
4. 评分应体现公平性，不应以专业设计师的标准要求参赛者
5. 提供具体、可操作的改进建议

输出格式：
请严格按照JSON格式输出，包含以下字段：
{
  "overall_score": 总分（0-100）,
  "creativity_score": 创意评分（0-20）,
  "visual_effect_score": 视觉效果评分（0-20）,
  "content_presentation_score": 内容呈现评分（0-20）,
  "originality_score": 原创性评分（0-20）,
  "documentation_score": 文档完整性评分（0-20）,
  "review_summary": "评审总结（200-500字）",
  "strengths": ["亮点1", "亮点2", ...],
  "weaknesses": ["不足1", "不足2", ...],
  "improvement_suggestions": ["建议1", "建议2", ...]
}"""

    def _build_review_prompt(
        self,
        ppt_metadata: Dict[str, Any],
        compliance_check: Dict[str, Any],
        work_description: str,
        readme_content: Optional[str],
        visual_effect_score: float,
        ppt_quality_evidence: Optional[Dict] = None,
    ) -> str:
        """构建PPT评审提示词（包含提取的文本内容）"""

        compliance_status = "符合所有硬性要求" if compliance_check["all_valid"] else "部分硬性要求不合规"

        # 构建PPT文本内容摘要（前10页）
        ppt_content_summary = ""
        if ppt_metadata.get("slide_contents"):
            ppt_content_summary = "\n【PPT文本内容摘要】（前10页）：\n"
            for slide in ppt_metadata["slide_contents"][:10]:
                ppt_content_summary += f"\n第{slide['slide_number']}页：\n"
                if slide['title']:
                    ppt_content_summary += f"标题：{slide['title']}\n"
                if slide['content']:
                    content_text = "\n".join(slide['content'][:3])
                    if len(content_text) > 200:
                        content_text = content_text[:200] + "..."
                    ppt_content_summary += f"内容：{content_text}\n"

        # 构建静态分析报告
        analysis_section = ""
        if ppt_quality_evidence:
            analysis_section = f"""
【静态分析参考】（系统自动检测，仅供参考）：

--- 视觉效果指标 ---
- 幻灯片数量：{ppt_quality_evidence.get('slide_count', 0)}页
- 有图片的页面比例：{ppt_quality_evidence.get('image_ratio', 0):.0%}
- 是否包含图表：{'是' if ppt_quality_evidence.get('has_chart') else '否'}
- 字体数量：{ppt_quality_evidence.get('font_count', 0)}种
- 字体一致性：{ppt_quality_evidence.get('font_consistency', '未知')}
- 建议分数参考：{ppt_quality_evidence.get('visual_score_hint', 8)}/20

--- 内容呈现指标 ---
- 平均每页文字量：{ppt_quality_evidence.get('avg_text_per_slide', 0)}字
- 文字密度：{ppt_quality_evidence.get('text_density_level', '未知')}
- 是否有结构（目录等）：{'有' if ppt_quality_evidence.get('has_structure') else '未检测到'}
- 建议分数参考：{ppt_quality_evidence.get('content_score_hint', 8)}/20

说明：以上"建议分数参考"是系统基于静态分析计算的参考值。请以实际PPT内容为主要依据进行评分。如果你的判断与参考值偏差较大（>5分），请在评审意见中说明原因。
"""

        prompt = f"""请评审以下演示文稿作品：

【PPT元数据】：
- 文件路径：{ppt_metadata['ppt_path']}
- 文件格式：{ppt_metadata['format']}
- 幻灯片数量：{ppt_metadata['slide_count']}页
- 幻灯片比例：{ppt_metadata['aspect_ratio']}
- 幻灯片尺寸：{ppt_metadata.get('slide_dimensions', '未知')}
- 文件大小：{ppt_metadata['file_size'] / (1024*1024):.2f}MB
- 密码保护：{'有' if ppt_metadata['has_password'] else '无'}
{ppt_content_summary}

【硬性要求合规性】：
状态：{compliance_status}
检查详情：
- 幻灯片数量：{'✓合规' if compliance_check['slide_count_valid'] else '✗不合规'}（{compliance_check['slide_count_message']}）
- 幻灯片比例：{'✓合规' if compliance_check['ratio_valid'] else '✗不合规'}（{compliance_check['ratio_message']}）
- 文件格式：{'✓合规' if compliance_check['format_valid'] else '✗不合规'}（{compliance_check['format_message']}）
- 密码保护：{'✓合规' if compliance_check['password_valid'] else '✗不合规'}（{compliance_check['password_message']}）
- 文件大小：{'✓合规' if compliance_check['size_valid'] else '✗不合规'}（{compliance_check['size_message']}）
- 说明文档：{'✓合规' if compliance_check['readme_exists'] else '✗不合规'}（{compliance_check['readme_message']}）
{analysis_section}
【作品说明】：
{work_description}

【README文档】：
{readme_content if readme_content else '未提供README文档'}

请根据评审标准，对演示文稿作品进行全面评价。注意参赛者为非计算机专业学生，请以鼓励为主，肯定其努力和亮点。请以JSON格式输出评审结果。"""

        return prompt

    def _parse_llm_output(
        self,
        llm_output: str,
        ppt_metadata: Dict[str, Any],
        compliance_check: Dict[str, Any],
        visual_effect_score: float
    ) -> PPTReviewOutput:
        """解析LLM输出为结构化评审结果"""
        try:
            result_dict = extract_json_from_llm_output(llm_output)

            # 构建PPTReviewOutput对象
            review_output = PPTReviewOutput(
                overall_score=result_dict.get("overall_score", visual_effect_score + 50),
                creativity_score=result_dict.get("creativity_score", 12),
                visual_effect_score=result_dict.get("visual_effect_score", visual_effect_score),
                content_presentation_score=result_dict.get("content_presentation_score", 12),
                originality_score=result_dict.get("originality_score", 12),
                documentation_score=result_dict.get("documentation_score", 8),
                compliance_check=compliance_check,
                review_summary=result_dict.get("review_summary", ""),
                strengths=result_dict.get("strengths", []),
                weaknesses=result_dict.get("weaknesses", []),
                improvement_suggestions=result_dict.get("improvement_suggestions", []),
                ppt_metadata=ppt_metadata
            )

            return review_output

        except json.JSONDecodeError as e:
            logger.error(f"LLM输出JSON解析失败：{e}")

            # 返回默认结果
            return PPTReviewOutput(
                overall_score=visual_effect_score + 50,
                creativity_score=12,
                visual_effect_score=visual_effect_score,
                content_presentation_score=12,
                originality_score=12,
                documentation_score=8,
                compliance_check=compliance_check,
                review_summary="评审结果解析失败",
                strengths=["评审解析失败"],
                weaknesses=["LLM输出格式错误"],
                improvement_suggestions=["请检查LLM输出格式"],
                ppt_metadata=ppt_metadata
            )


# 使用示例
if __name__ == "__main__":
    import os

    os.environ["DEEPSEEK_API_KEY"] = "your-deepseek-api-key"

    agent = PPTReviewerAgent()

    # 测试PPT路径（需要提供实际PPT）
    test_ppt = "test_ppt/demo.pptx"
    readme = """
    # 作品说明
    本演示文稿讲解人工智能基础知识。
    内容包括：AI定义、发展历程、应用场景等。
    共15页，采用16:9比例，原创设计。
    """

    if os.path.exists(test_ppt):
        print("="*80)
        print("PPTReviewerAgent完整测试")
        print("="*80)

        result = agent.review_ppt(
            ppt_path=test_ppt,
            work_description="人工智能基础演示文稿",
            readme_content=readme
        )

        print(f"\n总分：{result.overall_score}")
        print(f"创意：{result.creativity_score}")
        print(f"视觉效果：{result.visual_effect_score}")
        print(f"内容呈现：{result.content_presentation_score}")
        print(f"原创性：{result.originality_score}")

        print("="*80)
    else:
        print(f"测试PPT不存在：{test_ppt}")