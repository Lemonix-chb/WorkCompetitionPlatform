# File utilities - Extract and analyze submitted work files
import os
import zipfile
import tempfile
import shutil
import logging
from typing import Optional
from collections import Counter

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not installed, PPT extraction will be limited")

logger = logging.getLogger(__name__)

CODE_EXTENSIONS = {".java", ".py", ".js", ".ts", ".html", ".css", ".sql",
                   ".cpp", ".c", ".go", ".rs", ".jsx", ".tsx", ".vue", ".xml"}
DOC_EXTENSIONS = {".md", ".txt", ".rst", ".doc", ".docx"}
CONFIG_FILES = {"pom.xml", "package.json", "requirements.txt", "build.gradle",
                "Makefile", "Dockerfile", ".gitignore", "README.md", "README.txt",
                "README", "docker-compose.yml", "tsconfig.json", "vite.config.js"}
PPT_EXTENSIONS = {".pptx", ".ppt", ".pdf", ".key"}
VIDEO_EXTENSIONS = {".mp4", ".avi", ".mov", ".mkv", ".flv", ".wmv"}
IGNORED_DIRS = {"node_modules", "__pycache__", ".git", ".idea", "target",
                "dist", "build", ".vscode", "venv", ".venv", "env"}

MAX_LINES_PER_FILE = 200
MAX_TOTAL_LINES = 3000


def extract_and_analyze(file_path: str, work_type: str) -> dict:
    """
    提取并分析提交的作品文件

    Returns:
        dict: {
            file_manifest: str,    # 格式化的文件清单
            file_contents: str,    # 格式化的文件内容
            stats: dict,           # 统计信息
            duplicate_rate: None,  # 预留
            code_quality_score: None  # 预留
        }
    """
    logger.info(f"Extracting files: path={file_path}, type={work_type}")

    result = {
        "file_manifest": "（无文件）",
        "file_contents": "（无文件内容）",
        "stats": {},
        "duplicate_rate": None,
        "code_quality_score": None
    }

    if not file_path or not os.path.exists(file_path):
        logger.warning(f"File path does not exist: {file_path}")
        return result

    # 解压或直接读取目录
    work_dir = _get_work_directory(file_path)
    if not work_dir:
        logger.error(f"Cannot access file: {file_path}")
        return result

    try:
        # 收集文件信息
        all_files = _collect_files(work_dir)
        if not all_files:
            result["file_manifest"] = "（目录为空）"
            return result

        # 格式化文件清单
        result["file_manifest"] = _format_manifest(all_files)

        # 按类型筛选并读取内容
        contents, stats = _read_contents(all_files, work_type)
        result["file_contents"] = _format_contents(contents)
        result["stats"] = stats

    finally:
        # 清理临时目录
        if work_dir != file_path and work_dir.startswith(tempfile.gettempdir()):
            shutil.rmtree(work_dir, ignore_errors=True)

    return result


def _get_work_directory(file_path: str) -> Optional[str]:
    """获取工作目录（如果是ZIP则解压）"""
    if os.path.isdir(file_path):
        return file_path

    if file_path.lower().endswith(".zip"):
        return _unzip(file_path)

    # 单个文件，返回其所在目录
    parent = os.path.dirname(file_path)
    return parent if os.path.isdir(parent) else None


def _unzip(zip_path: str) -> Optional[str]:
    """解压ZIP文件到临时目录"""
    try:
        temp_dir = tempfile.mkdtemp(prefix="review_")
        with zipfile.ZipFile(zip_path, "r") as zf:
            zf.extractall(temp_dir)

        # 如果ZIP内只有一个顶层目录，进入该目录
        entries = os.listdir(temp_dir)
        if len(entries) == 1 and os.path.isdir(os.path.join(temp_dir, entries[0])):
            return os.path.join(temp_dir, entries[0])

        return temp_dir
    except Exception as e:
        logger.error(f"Unzip failed: {e}")
        return None


def _collect_files(directory: str) -> list[dict]:
    """递归收集目录中的所有文件"""
    files = []
    for root, dirs, filenames in os.walk(directory):
        # 过滤忽略的目录
        dirs[:] = [d for d in dirs if d not in IGNORED_DIRS]

        for name in filenames:
            full_path = os.path.join(root, name)
            rel_path = os.path.relpath(full_path, directory)
            ext = os.path.splitext(name)[1].lower()
            size = os.path.getsize(full_path)

            files.append({
                "name": name,
                "path": rel_path.replace("\\", "/"),
                "ext": ext,
                "size": size,
                "full_path": full_path
            })

    return sorted(files, key=lambda f: f["path"])


def _read_contents(files: list[dict], work_type: str) -> tuple[dict, dict]:
    """按作品类型筛选文件并读取内容"""
    contents = {}
    stats = {"code_file_count": 0, "total_lines": 0, "languages": [], "ppt_slides": 0}
    lang_counter = Counter()
    total_lines = 0

    for f in files:
        ext = f["ext"]
        name = f["name"]
        should_read = False

        if work_type == "CODE":
            if ext in CODE_EXTENSIONS or name in CONFIG_FILES or ext in DOC_EXTENSIONS:
                should_read = True
        elif work_type == "PPT":
            # 对于PPT作品，提取PPTX内容和文档文件
            if ext in PPT_EXTENSIONS:
                # 特殊处理PPTX文件
                if ext == ".pptx" and PPTX_AVAILABLE:
                    ppt_content = _extract_ppt_content(f["full_path"])
                    if ppt_content:
                        contents[f["path"]] = ppt_content
                        slides_count = ppt_content.count("[幻灯片")
                        stats["ppt_slides"] = slides_count
                        total_lines += slides_count
                continue
            elif ext in DOC_EXTENSIONS or name in CONFIG_FILES:
                should_read = True
        elif work_type == "VIDEO":
            if ext in DOC_EXTENSIONS or name in CONFIG_FILES:
                should_read = True

        if not should_read:
            continue

        if total_lines >= MAX_TOTAL_LINES:
            break

        content = _read_file(f["full_path"], MAX_LINES_PER_FILE)
        if content:
            contents[f["path"]] = content
            lines = content.count("\n") + 1
            total_lines += lines

            if ext in CODE_EXTENSIONS:
                stats["code_file_count"] += 1
                lang_name = _ext_to_language(ext)
                lang_counter[lang_name] += 1

    stats["total_lines"] = total_lines
    stats["languages"] = ", ".join(f"{lang}({cnt})" for lang, cnt in lang_counter.most_common(5))
    return contents, stats


def _read_file(path: str, max_lines: int) -> Optional[str]:
    """读取文件内容，带行数限制"""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            lines = []
            for i, line in enumerate(f):
                if i >= max_lines:
                    lines.append(f"... (truncated, showing first {max_lines} lines)")
                    break
                lines.append(line.rstrip())
            return "\n".join(lines)
    except Exception as e:
        logger.debug(f"Cannot read file {path}: {e}")
        return None


def _extract_ppt_content(pptx_path: str) -> Optional[str]:
    """提取PPT幻灯片的文本内容"""
    if not PPTX_AVAILABLE:
        logger.warning(f"python-pptx not available, cannot extract PPT content")
        return None

    try:
        prs = Presentation(pptx_path)
        logger.info(f"Loaded PPTX: {pptx_path}, slides: {len(prs.slides)}")

        slides_content = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            if slide_texts:
                slide_content = f"[幻灯片 {slide_num}]\n" + "\n".join(slide_texts)
                slides_content.append(slide_content)

        if slides_content:
            full_content = "\n\n".join(slides_content)
            logger.info(f"Extracted {len(slides_content)} slides from PPT")
            return full_content
        else:
            logger.warning(f"PPT has no text content")
            return None

    except Exception as e:
        logger.error(f"Failed to extract PPT content: {e}")
        return None


def _format_manifest(files: list[dict]) -> str:
    """格式化文件清单"""
    lines = []
    for f in files:
        size_kb = f["size"] / 1024
        if size_kb > 1024:
            size_str = f"{size_kb / 1024:.1f}MB"
        else:
            size_str = f"{size_kb:.1f}KB"
        lines.append(f"- {f['path']} ({size_str})")
    return "\n".join(lines)


def _format_contents(contents: dict[str, str]) -> str:
    """格式化文件内容用于提示词"""
    if not contents:
        return "（无可读取的文件内容）"

    parts = []
    for path, content in contents.items():
        parts.append(f"### {path}\n```\n{content}\n```")
    return "\n\n".join(parts)


def _ext_to_language(ext: str) -> str:
    """文件扩展名映射到语言名称"""
    mapping = {
        ".java": "Java", ".py": "Python", ".js": "JavaScript", ".ts": "TypeScript",
        ".html": "HTML", ".css": "CSS", ".sql": "SQL", ".cpp": "C++",
        ".c": "C", ".go": "Go", ".rs": "Rust", ".jsx": "React JSX",
        ".tsx": "React TSX", ".vue": "Vue", ".xml": "XML"
    }
    return mapping.get(ext, ext)
